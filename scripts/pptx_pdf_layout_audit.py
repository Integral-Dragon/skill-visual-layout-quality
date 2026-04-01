#!/usr/bin/env python3

"""Audit PPTX and PDF files for obvious layout-risk preflight issues.

This script is intentionally lightweight.

- `.pptx` support is stronger: it checks slide aspect ratio, scans text frames,
  estimates obvious overflow risk, and surfaces font-fragility patterns.
- `.pdf` support is metadata-level only via `pdfinfo`: page count, page size,
  and aspect-ratio drift. It does not inspect rendered text frames inside the
  PDF.
- if a PPTX run does not carry an explicit font size, the script falls back to
  a simple 18pt/20pt heuristic based on shape height. That is good enough for
  preflight warnings, but theme-level sizes can still make the estimate wrong.

Final judgment still belongs to rendered thumbnails or page images.
"""

from __future__ import annotations

import argparse
import re
import subprocess
import sys
from dataclasses import dataclass
from pathlib import Path

EMU_PER_PT = 12700
SAFE_PORTABLE_FONTS = {
    "arial",
    "aptos",
    "aptos display",
    "calibri",
    "cambria",
    "courier new",
    "georgia",
    "helvetica",
    "palatino",
    "tahoma",
    "times new roman",
    "trebuchet ms",
    "verdana",
}


def emu_to_pt(value: int | float | None, fallback: float = 0.0) -> float:
    if value is None:
        return fallback
    return float(value) / EMU_PER_PT


def normalize_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def is_all_caps(text: str) -> bool:
    letters = [char for char in text if char.isalpha()]
    return bool(letters) and all(char.upper() == char for char in letters)


def is_cjk(char: str) -> bool:
    code = ord(char)
    return (
        0x3400 <= code <= 0x4DBF
        or 0x4E00 <= code <= 0x9FFF
        or 0x3040 <= code <= 0x30FF
        or 0xAC00 <= code <= 0xD7AF
        or 0xF900 <= code <= 0xFAFF
        or 0xFF66 <= code <= 0xFF9D
    )


def is_monospace(font_name: str | None) -> bool:
    if not font_name:
        return False
    normalized = font_name.lower()
    keywords = (
        "mono",
        "courier",
        "consolas",
        "menlo",
        "monaco",
        "code",
        "jetbrains mono",
        "ibm plex mono",
        "sfmono",
    )
    return any(keyword in normalized for keyword in keywords)


def latin_ratio(text: str, font_name: str | None, is_bold: bool) -> float:
    if is_monospace(font_name):
        return 0.60
    caps = is_all_caps(text)
    if caps and is_bold:
        return 0.67
    if caps:
        return 0.62
    if is_bold:
        return 0.60
    return 0.55


def estimate_line_width(text: str, font_size_pt: float, font_name: str | None, is_bold: bool) -> float:
    ratio = latin_ratio(text, font_name, is_bold)
    width = 0.0
    for char in text:
        if char.isspace():
            width += font_size_pt * (0.60 if is_monospace(font_name) else 0.33)
        elif is_cjk(char):
            width += font_size_pt * 1.0
        else:
            width += font_size_pt * ratio
    return width


def aspect_label(width: float, height: float) -> str:
    ratio = width / height if height else 0.0
    if abs(ratio - (16 / 9)) < 0.03:
        return "16:9"
    if abs(ratio - (4 / 3)) < 0.03:
        return "4:3"
    return f"custom ({ratio:.2f}:1)"


def iter_shapes(shapes, group_type):
    for shape in shapes:
        yield shape
        if shape.shape_type == group_type:
            yield from iter_shapes(shape.shapes, group_type)


@dataclass
class TextFrameAudit:
    slide_number: int
    shape_name: str
    text: str
    available_width_pt: float
    available_height_pt: float
    estimated_width_pt: float
    estimated_height_pt: float
    font_names: set[str]
    missing_font_runs: int
    line_count: int


def audit_text_frame(slide_number: int, shape) -> TextFrameAudit | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    text_frame = shape.text_frame
    text = normalize_text(text_frame.text)
    if not text:
        return None

    font_names: set[str] = set()
    missing_font_runs = 0
    max_line_width = 0.0
    line_count = 0
    total_height = 0.0

    for paragraph in text_frame.paragraphs:
        para_text = normalize_text("".join(run.text for run in paragraph.runs) or paragraph.text or "")
        if not para_text:
            continue
        line_count += 1
        runs = [run for run in paragraph.runs if run.text]
        run_fonts = {run.font.name for run in runs if run.font.name}
        if run_fonts:
            font_names.update(run_fonts)
        missing_font_runs += sum(1 for run in runs if not run.font.name)
        font_name = next(iter(run_fonts), None)
        font_size = next((float(run.font.size.pt) for run in runs if run.font.size), 0.0)
        if not font_size:
            font_size = 18.0 if shape.height < 800000 else 20.0
        is_bold = any(bool(run.font.bold) for run in runs)
        max_line_width = max(max_line_width, estimate_line_width(para_text, font_size, font_name, is_bold))
        total_height += font_size * 1.25

    if line_count == 0:
        return None

    available_width = max(
        emu_to_pt(shape.width) - emu_to_pt(text_frame.margin_left) - emu_to_pt(text_frame.margin_right),
        0.0,
    )
    available_height = max(
        emu_to_pt(shape.height) - emu_to_pt(text_frame.margin_top) - emu_to_pt(text_frame.margin_bottom),
        0.0,
    )

    return TextFrameAudit(
        slide_number=slide_number,
        shape_name=shape.name or "<unnamed>",
        text=text,
        available_width_pt=available_width,
        available_height_pt=available_height,
        estimated_width_pt=max_line_width,
        estimated_height_pt=total_height,
        font_names=font_names,
        missing_font_runs=missing_font_runs,
        line_count=line_count,
    )


def audit_pptx(path: Path) -> int:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        print(path)
        print("  - python-pptx is not installed; install dependencies with `pip install -r requirements.txt`")
        return 1

    prs = Presentation(path)
    width_pt = emu_to_pt(prs.slide_width)
    height_pt = emu_to_pt(prs.slide_height)
    ratio_name = aspect_label(width_pt, height_pt)

    problems: list[str] = []
    notes: list[str] = []
    if ratio_name != "16:9":
        problems.append(
            f"aspect risk: deck is {ratio_name} ({width_pt:.0f}pt x {height_pt:.0f}pt); re-check target surface before export"
        )

    used_fonts: set[str] = set()
    risky_font_names: set[str] = set()
    inherited_font_frames = 0

    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in iter_shapes(slide.shapes, MSO_SHAPE_TYPE.GROUP):
            audit = audit_text_frame(slide_number, shape)
            if audit is None:
                continue
            label = audit.text[:72]

            if audit.estimated_width_pt > audit.available_width_pt:
                problems.append(
                    f"overflow risk: slide {audit.slide_number} '{audit.shape_name}' est={audit.estimated_width_pt:.0f}pt "
                    f"available={audit.available_width_pt:.0f}pt text='{label}'"
                )

            if audit.estimated_height_pt > audit.available_height_pt:
                problems.append(
                    f"height risk: slide {audit.slide_number} '{audit.shape_name}' est={audit.estimated_height_pt:.0f}pt "
                    f"available={audit.available_height_pt:.0f}pt text='{label}'"
                )

            if audit.available_width_pt < 96 and len(audit.text) > 12:
                problems.append(
                    f"dense-frame risk: slide {audit.slide_number} '{audit.shape_name}' is narrow ({audit.available_width_pt:.0f}pt) for text='{label}'"
                )

            if audit.missing_font_runs and audit.font_names:
                problems.append(
                    f"font risk: slide {audit.slide_number} '{audit.shape_name}' mixes explicit and inherited font runs"
                )
            elif audit.missing_font_runs and not audit.font_names:
                inherited_font_frames += 1

            if len(audit.font_names) > 1:
                problems.append(
                    f"font drift risk: slide {audit.slide_number} '{audit.shape_name}' mixes fonts {sorted(audit.font_names)}"
                )

            used_fonts.update(audit.font_names)
            risky_font_names.update(
                font for font in audit.font_names if font.lower() not in SAFE_PORTABLE_FONTS
            )

    if risky_font_names:
        problems.append(f"font portability risk: deck uses non-portable fonts {sorted(risky_font_names)}")
    if inherited_font_frames:
        notes.append(
            f"theme-font dependency: {inherited_font_frames} text frame(s) rely on inherited/theme fonts; verify export environment"
        )

    print(path)
    print(f"  deck: {len(prs.slides)} slide(s), aspect {ratio_name}")
    if used_fonts:
        print(f"  fonts seen: {', '.join(sorted(used_fonts))}")
    if problems:
        for problem in problems:
            print(f"  - {problem}")
        for note in notes:
            print(f"  note: {note}")
        return 1

    print("  OK: no obvious aspect, text-frame, or font-fragility risks found")
    for note in notes:
        print(f"  note: {note}")
    return 0


def parse_pdfinfo_page_size(output: str) -> tuple[float, float] | None:
    match = re.search(r"Page size:\s+([\d.]+)\s+x\s+([\d.]+)\s+pts", output)
    if not match:
        return None
    return float(match.group(1)), float(match.group(2))


def parse_pdfinfo_pages(output: str) -> int | None:
    match = re.search(r"Pages:\s+(\d+)", output)
    return int(match.group(1)) if match else None


def audit_pdf(path: Path) -> int:
    try:
        result = subprocess.run(
            ["pdfinfo", str(path)],
            capture_output=True,
            text=True,
            check=True,
        )
    except FileNotFoundError:
        print(path)
        print("  - pdfinfo is not installed; PDF audit is metadata-limited and unavailable here")
        return 1
    except subprocess.CalledProcessError as exc:
        print(path)
        stderr = exc.stderr.strip() or "pdfinfo failed"
        print(f"  - unable to inspect PDF metadata: {stderr}")
        return 1

    output = result.stdout
    pages = parse_pdfinfo_pages(output)
    size = parse_pdfinfo_page_size(output)
    problems: list[str] = []
    notes: list[str] = []

    print(path)
    if pages is not None:
        print(f"  pages: {pages}")
    if size is None:
        print("  - page-size metadata unavailable; rendered inspection still required")
        return 1

    width_pt, height_pt = size
    ratio_name = aspect_label(width_pt, height_pt)
    print(f"  first-page size: {width_pt:.0f}pt x {height_pt:.0f}pt ({ratio_name})")

    if ratio_name == "4:3":
        problems.append(
            "aspect risk: PDF is 4:3; if this came from slides, confirm the deck was not exported with the wrong page size"
        )
    elif ratio_name.startswith("custom"):
        notes.append(
            "custom page ratio: this may be intentional for document PDFs, but re-check if the source was supposed to be widescreen slides"
        )

    notes.append(
        "pdf preflight only: inspect rendered pages for overflow, clipping, font substitution, and footer collisions"
    )

    if problems:
        for problem in problems:
            print(f"  - {problem}")
        for note in notes:
            print(f"  note: {note}")
        return 1

    print("  OK: no obvious metadata-level aspect risks found")
    for note in notes:
        print(f"  note: {note}")
    return 0


def audit_path(path: Path) -> int:
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return audit_pptx(path)
    if suffix == ".pdf":
        return audit_pdf(path)
    print(f"{path}\n  - unsupported file type; use .pptx or .pdf", file=sys.stderr)
    return 1


def main() -> int:
    parser = argparse.ArgumentParser(description="Audit PPTX/PDF files for likely layout risks.")
    parser.add_argument("paths", nargs="+", help="PPTX or PDF files to inspect")
    args = parser.parse_args()

    exit_code = 0
    for raw_path in args.paths:
        path = Path(raw_path)
        if not path.exists():
            print(f"{path}\n  - missing file", file=sys.stderr)
            exit_code = 1
            continue
        exit_code = max(exit_code, audit_path(path))
    return exit_code


if __name__ == "__main__":
    raise SystemExit(main())
