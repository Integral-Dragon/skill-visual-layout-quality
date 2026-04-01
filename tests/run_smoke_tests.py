#!/usr/bin/env python3

from __future__ import annotations

import subprocess
import sys
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
SVG_AUDIT = ROOT / "scripts" / "svg_layout_audit.py"
PPTX_AUDIT = ROOT / "scripts" / "pptx_pdf_layout_audit.py"
SVG_FIXTURES = ROOT / "tests" / "fixtures" / "svg"


def run_and_expect(path: Path, expected_exit: int) -> None:
    script = SVG_AUDIT if path.suffix == ".svg" else PPTX_AUDIT
    result = subprocess.run(
        [sys.executable, str(script), str(path)],
        capture_output=True,
        text=True,
    )
    if result.returncode != expected_exit:
        raise SystemExit(
            f"{path.name}: expected exit {expected_exit}, got {result.returncode}\n"
            f"stdout:\n{result.stdout}\n\nstderr:\n{result.stderr}"
        )


def build_pptx(path: Path, *, narrow: bool) -> None:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise SystemExit(
            "python-pptx is required for PPTX smoke tests. Install dependencies with "
            "`pip install -r requirements.txt`."
        ) from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(1.4 if narrow else 4.8), Inches(0.6))
    text_frame = textbox.text_frame
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "This heading is intentionally too long for the frame" if narrow else "Stable headline"
    run.font.name = "Arial"
    run.font.size = Pt(20)
    prs.save(path)


def main() -> int:
    run_and_expect(SVG_FIXTURES / "good.svg", 0)
    run_and_expect(SVG_FIXTURES / "bad.svg", 1)

    with tempfile.TemporaryDirectory() as tmpdir:
        tmp = Path(tmpdir)
        good_pptx = tmp / "good.pptx"
        bad_pptx = tmp / "bad.pptx"
        build_pptx(good_pptx, narrow=False)
        build_pptx(bad_pptx, narrow=True)
        run_and_expect(good_pptx, 0)
        run_and_expect(bad_pptx, 1)

    print("Smoke tests passed")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
